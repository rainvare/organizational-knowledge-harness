"""parsers/pptx_parser.py"""


class PPTXParser:
    def __init__(self, api_key: str = ""):
        self.api_key = api_key

    def parse(self, path: str) -> dict:
        try:
            from pptx import Presentation
            prs = Presentation(path)
            texts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        texts.append(shape.text.strip())
            return {"text": "\n".join(texts), "modality": "pptx"}
        except ImportError:
            return {"text": "", "error": "python-pptx not installed"}
        except Exception as e:
            return {"text": "", "error": str(e)}
